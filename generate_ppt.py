from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

def create_presentation(output_path):
    prs = Presentation()

    # 1. Title Slide
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "Address Geocoding & Normalization System"
    subtitle.text = "Comprehensive Project Overview & Operation Guide\nPrepared for: Kishor Wakchaure"

    # 2. Project Objective
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Project Objective"
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "Build a robust company address geocoding registry that:"
    p = tf.add_paragraph()
    p.text = "• Standardizes company names and addresses."
    p = tf.add_paragraph()
    p.text = "• Minimizes API costs through intelligent caching."
    p = tf.add_paragraph()
    p.text = "• Provides a collaborative workflow using Google Sheets."
    p = tf.add_paragraph()
    p.text = "• Enables analysts via CLI and Web interfaces."

    # 3. Key Components
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Key Components"
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "The system consists of 4 main modules:"
    p = tf.add_paragraph()
    p.text = "• Normalization Engine: Cleans and standardizes names."
    p = tf.add_paragraph()
    p.text = "• Geocoding Service: Integration with Google Maps API."
    p = tf.add_paragraph()
    p.text = "• Multi-Tier Storage: Google Sheets + SQLite + Memory Cache."
    p = tf.add_paragraph()
    p.text = "• Quality Assurance: Automated confidence scoring and review queue."

    # 4. Smart Normalization
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Smart Normalization"
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "Ensuring consistency across lookups:"
    p = tf.add_paragraph()
    p.text = "• Unicode Normalization: Handles special characters."
    p = tf.add_paragraph()
    p.text = "• Suffix Removal: Strips 'Ltd', 'LLC', 'Pvt Ltd' etc."
    p = tf.add_paragraph()
    p.text = "• Golden Mappings: Expands acronyms (e.g., TCS → Tata Consultancy Services)."
    p = tf.add_paragraph()
    p.text = "• Case Standardization: Everything converted to uppercase."

    # 5. Architecture & Workflow
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Architecture & Workflow"
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "Intelligent sequence to minimize costs:"
    p = tf.add_paragraph()
    p.text = "1. Cache Check: Instant hit for recent queries (FREE)."
    p = tf.add_paragraph()    
    p.text = "2. Registry Check: Check Google Sheets for existing records (FREE)."
    p = tf.add_paragraph()
    p.text = "3. Fuzzy Match: Search similar names in existing data (FREE)."
    p = tf.add_paragraph()
    p.text = "4. API Call: Only geocode if truly new (Costs API credit)."
    p = tf.add_paragraph()
    p.text = "5. Store & Sync: Local cache and Registry updated automatically."

    # 6. User Interfaces
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "User Interfaces"
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "System accessibility:"
    p = tf.add_paragraph()
    p.text = "• Web App (Streamlit): Visual, interactive maps, CSV batch upload."
    p = tf.add_paragraph()
    p.text = "• CLI Tool: High-speed lookups, statistics, and review queue."
    p = tf.add_paragraph()
    p.text = "• Google Sheets: Direct access to the address database."

    # 7. Deployment: Streamlit Cloud
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Cloud Deployment"
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "Deployed on Streamlit Community Cloud:"
    p = tf.add_paragraph()
    p.text = "• Connected directly to GitHub Repository."
    p = tf.add_paragraph()
    p.text = "• Live URL for team access."
    p = tf.add_paragraph()
    p.text = "• Dynamic Configuration: Users enter API keys in the app UI."
    p = tf.add_paragraph()
    p.text = "• Auto-deployment: Updates push automatically from GitHub."

    # 8. Summary of Achievements
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Summary of Work"
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "Completed Deliverables:"
    p = tf.add_paragraph()
    p.text = "✓ Full end-to-end Python backend."
    p = tf.add_paragraph()
    p.text = "✓ Modern Streamlit web interface."
    p = tf.add_paragraph()
    p.text = "✓ Google Sheets & API integration."
    p = tf.add_paragraph()
    p.text = "✓ Deployment on GitHub and Streamlit Cloud."
    p = tf.add_paragraph()
    p.text = "✓ Comprehensive documentation and setup guides."

    prs.save(output_path)
    print(f"Presentation saved to: {output_path}")

if __name__ == "__main__":
    create_presentation("Address_Geocoding_System_Overview.pptx")
