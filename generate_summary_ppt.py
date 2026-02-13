from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

def create_presentation():
    prs = Presentation()

    # Title Slide
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "Address Geocoding & Verification System"
    subtitle.text = "Intelligent Company Data Standardization\nProject Overview & Results"

    # Slide 1: The Problem
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "The Problem: Inconsistent Data"
    content = slide.placeholders[1].text_frame
    content.text = "Managing company address data often leads to several challenges:"
    p = content.add_paragraph()
    p.text = "• Manual entry results in inconsistent formats and typos."
    p.level = 1
    p = content.add_paragraph()
    p.text = "• Geocoding APIs (like Google Maps) are expensive if called repeatedly."
    p.level = 1
    p = content.add_paragraph()
    p.text = "• Verifying the accuracy of geocoded results against reality is time-consuming."
    p.level = 1
    p = content.add_paragraph()
    p.text = "• No central registry for teams to share already-found addresses."
    p.level = 1

    # Slide 2: The Solution
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "The Solution Architecture"
    content = slide.placeholders[1].text_frame
    content.text = "A multi-tier geocoding system designed for cost-efficiency and accuracy:"
    p = content.add_paragraph()
    p.text = "• Multi-Tier Caching: Checking Session, SQLite Local, and Global Google Sheets."
    p.level = 1
    p = content.add_paragraph()
    p.text = "• Smart Normalization: Standardizing company names before searching."
    p.level = 1
    p = content.add_paragraph()
    p.text = "• Global Sheets Registry: A shared database for team-wide address reuse."
    p.level = 1
    p = content.add_paragraph()
    p.text = "• Streamlit Interface: Professional web UI for easy interaction."
    p.level = 1

    # Slide 3: Fully Agentic Mode
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Fully Agentic Mode (Powered by Gemini AI)"
    content = slide.placeholders[1].text_frame
    content.text = "The system's most advanced feature for automated verification:"
    p = content.add_paragraph()
    p.text = "• AI Verification: LLM searches the web to find the company's official website."
    p.level = 1
    p = content.add_paragraph()
    p.text = "• Cross-Check: Compares the Google Maps result against the company's 'Contact Us' page."
    p.level = 1
    p = content.add_paragraph()
    p.text = "• Confidence Scoring: Assigns an AI confidence score and provides insights."
    p.level = 1
    p = content.add_paragraph()
    p.text = "• Source Transparency: Provides the direct URL used for AI verification."
    p.level = 1

    # Slide 4: Key Features
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "System Features"
    content = slide.placeholders[1].text_frame
    p = content.add_paragraph()
    p.text = "1. Individual Lookup: Precise geocoding with map visualization."
    p = content.add_paragraph()
    p.text = "2. Batch Processing: Upload CSV files to process thousands of records."
    p = content.add_paragraph()
    p.text = "3. Standardized Output: Aligning to custom 'STREET ADDRESS1', 'CITY' labels."
    p = content.add_paragraph()
    p.text = "4. Review Queue: Dedicated interface for manual QA of low-confidence results."
    p = content.add_paragraph()
    p.text = "5. Global Analytics: View storage stats and source distribution."

    # Slide 5: How it Works (Process)
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "How the App Works (Workflow)"
    content = slide.placeholders[1].text_frame
    p = content.add_paragraph()
    p.text = "User Input → Company Normalization →"
    p = content.add_paragraph()
    p.text = "→ Check Local Cache → Check Google Sheets Storage →"
    p = content.add_paragraph()
    p.text = "→ Call Google Maps Geocoding API (only if required) →"
    p = content.add_paragraph()
    p.text = "→ Optional: Agentic AI Verification (Gemini Web Search) →"
    p = content.add_paragraph()
    p.text = "→ Save Result Everywhere & Display to User."

    # Final Summary
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    title = slide.shapes.title
    title.text = "Next Steps & Future Scaling"
    subtitle = slide.placeholders[1]
    subtitle.text = "The system is now production-ready and fully scalable to over 500,000 records per sheet."

    # Save
    prs.save('Project_Overview_Address_Geocoding.pptx')
    print("Success: PPT generated as Project_Overview_Address_Geocoding.pptx")

if __name__ == "__main__":
    create_presentation()
